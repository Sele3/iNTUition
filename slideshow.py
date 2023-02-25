import os
import collections 
import collections.abc
from pptx import Presentation  
from pptx.util import Inches  


def generate_powerpoint() -> None:
    dst = "./frontend/public/media/"       # Images destination
    images = os.listdir(dst)    # Get their names in a list
    length = len(images)

    myPPT = Presentation()  

    for image in images:
        if '.png' not in image:
            continue

        # defining the path of the image file  
        imgPath = dst + image
        
        # creating the slide layout and adding a new slide  
        slideLayout = myPPT.slide_layouts[6]  
        mySlide = myPPT.slides.add_slide(slideLayout)  
        
        # specifying the values of the parameters for the add_picture() method  
        left = top = Inches(1)  
        # using the add_picture() method  
        myImage = mySlide.shapes.add_picture(imgPath, left, top)  
    
    # saving the PPT file  
    myPPT.save(dst + 'myPPT.pptx')  
